import sys
import os
import argparse

def convert_pdf_to_docx(input_path, output_path):
    print(f"Converting PDF to Word: {input_path} -> {output_path}")
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    print("PDF to Word conversion completed successfully.")

def convert_pdf_to_xlsx(input_path, output_path):
    print(f"Converting PDF to Excel: {input_path} -> {output_path}")
    import pdfplumber
    from openpyxl import Workbook
    
    wb = Workbook()
    wb.remove(wb.active) # Remove default sheet
    
    with pdfplumber.open(input_path) as pdf:
        for i, page in enumerate(pdf.pages):
            sheet = wb.create_sheet(title=f"Page {i+1}")
            tables = page.extract_tables()
            if tables:
                for t_idx, table in enumerate(tables):
                    if t_idx > 0:
                        sheet.append([]) # Empty row separator
                    for row in table:
                        sheet.append([str(cell) if cell is not None else "" for cell in row])
            else:
                # Fallback: extract text and append it
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        sheet.append([line])
                else:
                    sheet.append(["[Empty Page]"])
                    
    wb.save(output_path)
    print("PDF to Excel conversion completed successfully.")

def convert_pdf_to_pptx(input_path, output_path):
    print(f"Converting PDF to PowerPoint: {input_path} -> {output_path}")
    import pdfplumber
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    prs = Presentation()
    
    with pdfplumber.open(input_path) as pdf:
        for i, page in enumerate(pdf.pages):
            # Layout 6 is a blank slide layout
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Extract and place text
            text = page.extract_text()
            if text:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Title or marker
                p = tf.paragraphs[0]
                p.text = f"Slide {i+1}"
                p.font.size = Pt(20)
                p.font.bold = True
                
                # Append text lines
                for line in text.split("\n"):
                    p2 = tf.add_paragraph()
                    p2.text = line
                    p2.font.size = Pt(13)
            else:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = f"Slide {i+1} (No text extracted)"
                p.font.size = Pt(16)
                
    prs.save(output_path)
    print("PDF to PowerPoint conversion completed successfully.")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="OfficeFlex PDF Converter Helper")
    parser.add_argument("--input", required=True, help="Input PDF file path")
    parser.add_argument("--output", required=True, help="Output file path")
    parser.add_argument("--format", required=True, choices=["docx", "xlsx", "pptx"], help="Target format")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.input):
        print(f"Error: Input file {args.input} does not exist.", file=sys.stderr)
        sys.exit(1)
        
    # Ensure parent output directory exists
    out_dir = os.path.dirname(args.output)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
        
    try:
        if args.format == "docx":
            convert_pdf_to_docx(args.input, args.output)
        elif args.format == "xlsx":
            convert_pdf_to_xlsx(args.input, args.output)
        elif args.format == "pptx":
            convert_pdf_to_pptx(args.input, args.output)
        else:
            print(f"Error: Unsupported format {args.format}", file=sys.stderr)
            sys.exit(1)
            
    except Exception as e:
        print(f"Conversion failed with error: {str(e)}", file=sys.stderr)
        sys.exit(1)
